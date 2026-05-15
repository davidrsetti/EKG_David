"""
scripts/build_exec_pptx.py — Generate 2-slide CDTO executive deck for NEXUS.

Slide 1: What NEXUS is + the 10 live capabilities
Slide 2: Competitive moat + roadmap
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, datetime

# ── Brand colours ──────────────────────────────────────────────────────────────
GSK_ORANGE  = RGBColor(0xF3, 0x66, 0x33)
DARK        = RGBColor(0x1A, 0x1A, 0x1A)
MID         = RGBColor(0x44, 0x44, 0x44)
LIGHT_GREY  = RGBColor(0xF5, 0xF5, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x08, 0x91, 0xB2)
ACCENT_PURP = RGBColor(0x7C, 0x3A, 0xED)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, l, t, w, h,
                 font_size=12, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=10, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Platform Overview
# ══════════════════════════════════════════════════════════════════════════════

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Orange header bar ─────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.333, 1.05, GSK_ORANGE)

    add_text_box(slide, "NEXUS", 0.35, 0.08, 3, 0.45,
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "Enterprise Knowledge Graph Platform",
                 0.35, 0.52, 6, 0.38,
                 font_size=13, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
    add_text_box(slide, "CDTO Executive Briefing  ·  May 2026",
                 8.5, 0.6, 4.5, 0.35,
                 font_size=9, color=RGBColor(0xFF, 0xCC, 0xAA), align=PP_ALIGN.RIGHT)

    # ── Tagline ───────────────────────────────────────────────────────────────
    add_rect(slide, 0, 1.05, 13.333, 0.55, LIGHT_GREY)
    add_text_box(
        slide,
        "The only platform that connects business capabilities → applications → data assets"
        " → AI agents → security findings in one queryable knowledge graph with natural-language access.",
        0.35, 1.1, 12.6, 0.45,
        font_size=10.5, bold=False, color=MID, align=PP_ALIGN.LEFT,
    )

    # ── Section: 10 Live Capabilities (2-col cards) ───────────────────────────
    add_text_box(slide, "10 CAPABILITIES  ·  LIVE TODAY",
                 0.35, 1.72, 5, 0.28,
                 font_size=7.5, bold=True, color=GSK_ORANGE, align=PP_ALIGN.LEFT)

    caps = [
        ("💬", "Knowledge Graph Chat",
         "NL → SPARQL → synthesised answer in seconds"),
        ("📊", "Portfolio Intelligence",
         "Auto-scores every app on Gartner TIME model from live graph"),
        ("🏥", "SA Portfolio Health",
         "6-dim scan: gaps, tech debt, orphans, hotspots, data risk"),
        ("🧭", "Guided SA Advisor",
         "Graph-grounded 4-step SA interview → recommendations"),
        ("🗺️", "Architecture Diagrams",
         "7 diagram types auto-generated; draw.io / ArchiMate export"),
        ("💥", "Change Impact Radar",
         "6 parallel traversals: dependents, caps, data, agents, people"),
        ("🤖", "AI Agent Governance",
         "Registry + risk tiers + 0–100 governance score from graph"),
        ("📄", "ADR Generation",
         "MADR-format decision records stored as queryable graph nodes"),
        ("📋", "Data Query",
         "NL → SQL over Databricks Unity Catalog"),
        ("🔍", "Audit & Observability",
         "Immutable audit log · PII detection · guard block metrics"),
    ]

    col_l = [0.35, 6.85]
    row_t = 2.08
    row_h = 0.54

    for i, (icon, title, desc) in enumerate(caps):
        col = i % 2
        row = i // 2
        lx = col_l[col]
        ty = row_t + row * row_h

        # Card bg
        add_rect(slide, lx, ty, 6.25, row_h - 0.06,
                 LIGHT_GREY if col == 0 else RGBColor(0xFD, 0xF6, 0xF3))

        # Orange left accent
        add_rect(slide, lx, ty, 0.05, row_h - 0.06, GSK_ORANGE)

        # Icon + title
        add_text_box(slide, f"{icon}  {title}",
                     lx + 0.12, ty + 0.03, 6.0, 0.22,
                     font_size=9.5, bold=True, color=DARK)
        add_text_box(slide, desc,
                     lx + 0.12, ty + 0.24, 6.0, 0.22,
                     font_size=8.5, bold=False, color=MID)

    # ── Bottom bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 7.18, 13.333, 0.32, DARK)
    add_text_box(
        slide,
        "Stack: Python · FastAPI · Streamlit · StarDog RDF/SPARQL · Databricks SQL · OpenAI o3-mini / GPT-4o"
        "  ·  Role-based JWT auth  ·  9-stage safe query pipeline  ·  Immutable audit log",
        0.35, 7.21, 12.6, 0.26,
        font_size=7, color=RGBColor(0xBB, 0xBB, 0xBB),
    )

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Competitive Moat + Roadmap
# ══════════════════════════════════════════════════════════════════════════════

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Orange header bar ─────────────────────────────────────────────────────
    add_rect(slide, 0, 0, 13.333, 1.05, GSK_ORANGE)
    add_text_box(slide, "Why No Competitor Can Replicate This",
                 0.35, 0.14, 9, 0.45,
                 font_size=22, bold=True, color=WHITE)
    add_text_box(slide, "Competitive Differentiation  ·  Strategic Roadmap",
                 0.35, 0.62, 7, 0.32,
                 font_size=11, color=RGBColor(0xFF, 0xCC, 0xAA))

    # ── Left column: Competitor table ────────────────────────────────────────
    add_text_box(slide, "STRUCTURAL MOAT vs. INCUMBENTS",
                 0.35, 1.15, 7.2, 0.28,
                 font_size=7.5, bold=True, color=GSK_ORANGE)

    rows = [
        ("LeanIX",             "Portfolio roadmaps",          "Manual entry · no graph · no LLM",              "Live graph TIME scoring · NL query"),
        ("Ardoq / Bizzdesign", "ArchiMate modelling",         "No NL · no data governance layer",              "NL→SPARQL + graph-grounded SA advisor"),
        ("Collibra / Alation", "Data catalog + lineage",      "No EA/SA · no capability model",                "EA + data + AI agents in one graph"),
        ("ServiceNow CMDB",    "Configuration management",    "CI dependencies only · no semantic layer",      "Full blast-radius: caps, data, people, agents"),
        ("Microsoft Purview",  "Azure data governance",       "No SA advisor · no capability model",           "Cross-domain: EA + data + AI + findings"),
        ("Palantir AIP",       "Powerful ontology",           "$$$ · opinionated · no guided SA flow",         "Approachable NL-native · guided interview"),
    ]

    # Header row
    hx, hy, hw = 0.35, 1.48, [2.1, 1.85, 2.45, 2.7]
    headers = ["Competitor", "Their Strength", "Their Gap", "NEXUS Advantage"]
    hcols   = [hx, hx+hw[0], hx+hw[0]+hw[1], hx+hw[0]+hw[1]+hw[2]]

    add_rect(slide, hx-0.05, hy-0.02, sum(hw)+0.1, 0.3, DARK)
    for h, hcol in zip(headers, hcols):
        add_text_box(slide, h, hcol, hy, 2.6, 0.28,
                     font_size=8, bold=True, color=WHITE)

    row_colours = [LIGHT_GREY, WHITE]
    for ri, (comp, strength, gap, nexus) in enumerate(rows):
        ry = hy + 0.3 + ri * 0.38
        add_rect(slide, hx-0.05, ry-0.02, sum(hw)+0.1, 0.38,
                 row_colours[ri % 2])
        vals = [comp, strength, gap, nexus]
        fcols = [DARK, MID, RED, GREEN]
        for v, hcol, fc in zip(vals, hcols, fcols):
            add_text_box(slide, v, hcol, ry, 2.6, 0.36,
                         font_size=8, color=fc,
                         bold=(fc == DARK))

    # ── Right column: Roadmap ─────────────────────────────────────────────────
    add_rect(slide, 9.55, 1.15, 3.45, 6.1, LIGHT_GREY)
    add_rect(slide, 9.55, 1.15, 3.45, 0.32, DARK)
    add_text_box(slide, "NEXT ON ROADMAP",
                 9.65, 1.18, 3.2, 0.26,
                 font_size=8, bold=True, color=WHITE)

    roadmap = [
        (ACCENT_PURP, "D-4  Semantic Search",
         "Vector + graph hybrid. Eliminates 'no results' from name mismatches; "
         "injects entity hints into SPARQL generation."),
        (AMBER, "D-5  Executive Briefing Engine",
         "Weekly auto-generated HTML portfolio report: TIME distribution, "
         "top-5 invest/eliminate, capability heat map."),
        (GREEN, "D-6  Standards Conformance",
         "EA rules authored in-graph, evaluated against it. "
         "Live pass/fail conformance dashboard."),
        (ACCENT_BLUE, "E-1  ServiceNow CMDB Sync",
         "Two-way sync. NEXUS becomes the semantic intelligence layer "
         "above CMDB."),
    ]

    for ri, (col, title, body) in enumerate(roadmap):
        ry = 1.57 + ri * 1.32
        add_rect(slide, 9.55, ry, 0.07, 0.95, col)
        add_text_box(slide, title,
                     9.72, ry + 0.04, 3.15, 0.28,
                     font_size=9, bold=True, color=col)
        add_text_box(slide, body,
                     9.72, ry + 0.3, 3.15, 0.65,
                     font_size=8, color=MID, wrap=True)

    # ── Bottom: Governance callout ─────────────────────────────────────────────
    add_rect(slide, 0.3, 6.78, 9.1, 0.55, RGBColor(0xFD, 0xF6, 0xF3))
    add_rect(slide, 0.3, 6.78, 0.06, 0.55, GSK_ORANGE)
    add_text_box(
        slide,
        "Security by design: 9-stage safe query pipeline · Role-based JWT · SPARQL injection guard "
        "· PII auto-redaction · Immutable audit log · Randomised JWT secret per process",
        0.45, 6.82, 9.0, 0.46,
        font_size=8, color=MID,
    )

    # ── Bottom bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 7.18, 13.333, 0.32, DARK)
    add_text_box(
        slide,
        "NEXUS v2  ·  Confidential  ·  May 2026",
        0.35, 7.21, 12.6, 0.26,
        font_size=7.5, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.RIGHT,
    )

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide1(prs)
    build_slide2(prs)

    out = "/Users/drs58706/david/EKG_David/EKG_David/NEXUS_CDTO_Executive_Brief.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
